#!/usr/bin/env python3
"""
PowerPoint Presentation Generator
A complete, self-contained Python script to automatically generate professional PowerPoint presentations
using LLM content generation, real-time web search results, and image integration.

Author: AI Assistant
Date: 2024
"""

import os
import json
import re
import requests
import tempfile
import sys
import subprocess
from datetime import datetime
from typing import Dict, List, Optional

# --- Third-party library imports ---
# Ensure these are installed via pip
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
except ImportError:
    print("Error: python-pptx library not found. Install with: pip install python-pptx")
    sys.exit(1)

try:
    import google.generativeai as genai
except ImportError:
    print("Error: google-generativeai library not found. Install with: pip install google-generativeai")
    sys.exit(1)

try:
    from serpapi import GoogleSearch
except ImportError:
    print("Error: google-search-results library not found. Install with: pip install google-search-results")
    sys.exit(1)

try:
    from dotenv import load_dotenv
    load_dotenv()  # Load environment variables from .env file
except ImportError:
    print("Warning: python-dotenv not found. Install with: pip install python-dotenv")
    print("You can still set API keys as environment variables manually.")


class PowerPointGenerator:
    """
    Main class for generating PowerPoint presentations using LLM and web search.
    """

    def __init__(self):
        """Initializes the PowerPoint generator with API configurations."""
        self.gemini_api_key = os.getenv('GEMINI_API_KEY')
        self.serpapi_api_key = os.getenv('SERPAPI_API_KEY')
        
        # Initialize Gemini AI
        try:
            genai.configure(api_key=self.gemini_api_key)
            self.model = genai.GenerativeModel('gemini-2.5-flash')
        except Exception as e:
            # This is a safe place to handle configuration errors without crashing on import
            self.model = None

        # Define the basic slide structure
        self.slide_structure = {
            "title_slide": {"title": "", "subtitle": ""},
            "overview": {"title": "Overview", "content": []},
            "key_points": [
                {"title": "", "content": []},
                {"title": "", "content": []},
                {"title": "", "content": []},
                {"title": "", "content": []}
            ],
            "conclusion": {"title": "Conclusion", "content": []}
        }
    
    def _add_title_slide(self, prs: Presentation, content: Dict):
        """Adds the title slide to the presentation."""
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = content["title_slide"]["title"]
        subtitle.text = content["title_slide"]["subtitle"]

    def _add_content_slide(self, prs: Presentation, title: str, content_list: List[str], image_path: Optional[str] = None):
        """Adds a content slide with a title, bullet points, and an optional image."""
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = title
        
        text_frame = slide.placeholders[1].text_frame
        for point in content_list:
            p = text_frame.add_paragraph()
            p.text = point
            p.level = 0
            
        if image_path:
            try:
                left = Inches(6.5)
                top = Inches(1.5)
                height = Inches(4)
                slide.shapes.add_picture(image_path, left, top, height=height)
                print(f"🖼️ Image added to slide for: {title}")
            except Exception as e:
                print(f"❌ Could not add image to slide for: {title}. Error: {e}")
            finally:
                if image_path and os.path.exists(image_path):
                    os.unlink(image_path)
    
    def get_user_topic(self) -> str:
        """Prompts the user for presentation topic and returns the input."""
        print("\n" + "="*60)
        print("🎯 POWERPOINT PRESENTATION GENERATOR")
        print("="*60)
        print("\nEnter a topic for your presentation.")
        print("Examples:")
        print("  • The future of renewable energy")
        print("  • Advancements in quantum computing")
        print("  • The impact of AI on creative industries")
        print("  • Climate change mitigation strategies")
        print("  • Blockchain technology in healthcare")
        print("\nYour topic: ", end="")
        
        topic = input().strip()
        if not topic:
            print("Error: Topic cannot be empty. Please try again.")
            return self.get_user_topic()
        
        return topic
    
    def perform_web_search(self, topic: str) -> List[str]:
        """Performs a web search to get recent information on the topic."""
        print(f"\n🔍 Searching the web for recent information on: {topic}")
        
        try:
            search_params = {
                "q": f"{topic} latest news developments 2024",
                "api_key": self.serpapi_api_key,
                "engine": "google",
                "num": 5,
                "gl": "us",
                "hl": "en"
            }
            
            search = GoogleSearch(search_params)
            results = search.get_dict()
            
            snippets = []
            if "organic_results" in results:
                snippets = [result["snippet"] for result in results["organic_results"] if "snippet" in result][:5]
            
            if snippets:
                print(f"✅ Found {len(snippets)} relevant search results")
                return snippets
            else:
                print("⚠️  No search results found, proceeding with LLM generation only")
                return []
        except Exception as e:
            print(f"⚠️  Web search failed: {e}. Proceeding with LLM generation only...")
            return []

    def perform_image_search(self, query: str) -> Optional[str]:
        """Performs an image search and returns the URL of the first result."""
        print(f"🖼️ Searching for a relevant image for: {query}")
        try:
            search_params = {
                "q": query,
                "engine": "google_images",
                "api_key": self.serpapi_api_key,
                "ijn": "0" # Page number
            }
            search = GoogleSearch(search_params)
            results = search.get_dict()
            
            if "images_results" in results and len(results["images_results"]) > 0:
                first_image_url = results["images_results"][0].get("original")
                print("✅ Found an image.")
                return first_image_url
            else:
                print("⚠️ No image found.")
                return None
        except Exception as e:
            print(f"❌ Image search failed: {e}")
            return None

    def get_image_from_url(self, url: str) -> Optional[str]:
        """Downloads an image from a URL and saves it to a temporary file."""
        try:
            response = requests.get(url, stream=True, timeout=10)
            if response.status_code == 200:
                temp_file = tempfile.NamedTemporaryFile(delete=False, suffix=".jpg")
                with open(temp_file.name, 'wb') as f:
                    for chunk in response.iter_content(chunk_size=1024):
                        f.write(chunk)
                return temp_file.name
            else:
                return None
        except requests.exceptions.RequestException as e:
            print(f"❌ Failed to download image from {url}: {e}")
            return None
    
    def generate_llm_content(self, topic: str, web_results: List[str]) -> Dict:
        """Generates presentation content using Gemini AI."""
        print(f"\n🤖 Generating presentation content using AI...")
        
        web_context = ""
        if web_results:
            web_context = "\n\nRecent web search results to incorporate:\n" + "\n".join([f"- {snippet}" for snippet in web_results[:3]])
        
        prompt = f"""
You are an expert presentation creator. Create a comprehensive PowerPoint presentation on: "{topic}"

{web_context}

Generate the presentation content in the following EXACT JSON format. Do not include any additional text or explanations:

{{
    "title_slide": {{
        "title": "A compelling, professional title for the presentation",
        "subtitle": "{topic}"
    }},
    "overview": {{
        "title": "Overview",
        "content": [
            "Key point 1 about the topic",
            "Key point 2 about the topic", 
            "Key point 3 about the topic",
            "Key point 4 about the topic"
        ]
    }},
    "key_points": [
        {{
            "title": "Specific heading for slide 3 (e.g., 'Key Trends', 'Current State', 'Background')",
            "content": [
                "Detailed bullet point 1",
                "Detailed bullet point 2",
                "Detailed bullet point 3",
                "Detailed bullet point 4",
                "Detailed bullet point 5"
            ]
        }},
        {{
            "title": "Specific heading for slide 4 (e.g., 'Challenges', 'Issues', 'Problems')",
            "content": [
                "Detailed bullet point 1",
                "Detailed bullet point 2",
                "Detailed bullet point 3",
                "Detailed bullet point 4",
                "Detailed bullet point 5"
            ]
        }},
        {{
            "title": "Specific heading for slide 5 (e.g., 'Innovations', 'Solutions', 'Opportunities')",
            "content": [
                "Detailed bullet point 1",
                "Detailed bullet point 2",
                "Detailed bullet point 3",
                "Detailed bullet point 4",
                "Detailed bullet point 5"
            ]
        }},
        {{
            "title": "Specific heading for slide 6 (e.g., 'Impact', 'Future', 'Implications')",
            "content": [
                "Detailed bullet point 1",
                "Detailed bullet point 2",
                "Detailed bullet point 3",
                "Detailed bullet point 4",
                "Detailed bullet point 5"
            ]
        }}
    ],
    "conclusion": {{
        "title": "Conclusion",
        "content": [
            "Main takeaway 1",
            "Main takeaway 2", 
            "Main takeaway 3",
            "Concluding statement or call to action"
        ]
    }}
}}

IMPORTANT: 
- Ensure all content is relevant to "{topic}"
- Make bullet points informative and professional
- Keep titles concise and engaging
- Return ONLY valid JSON, no additional text
"""
        
        try:
            response = self.model.generate_content(prompt)
            content = response.text.strip()
            
            json_match = re.search(r'\{.*\}', content, re.DOTALL)
            if json_match:
                json_content = json_match.group()
                parsed_content = json.loads(json_content)
                print("✅ AI content generation successful")
                return parsed_content
            else:
                raise ValueError("No valid JSON found in response")
                
        except Exception as e:
            print(f"❌ AI content generation failed: {e}. Falling back to template content...")
            return self._get_fallback_content(topic)
    
    def _get_fallback_content(self, topic: str) -> Dict:
        """Generates fallback content if AI generation fails."""
        return {
            "title_slide": {
                "title": f"Comprehensive Analysis: {topic.title()}",
                "subtitle": topic
            },
            "overview": {
                "title": "Overview",
                "content": [
                    f"Introduction to {topic}",
                    "Current state and significance",
                    "Key areas of focus",
                    "Expected outcomes and insights"
                ]
            },
            "key_points": [
                {
                    "title": "Background & Context",
                    "content": [
                        "Historical development and evolution",
                        "Current market or industry status",
                        "Key stakeholders and participants",
                        "Regulatory and policy framework",
                        "Global trends and patterns"
                    ]
                },
                {
                    "title": "Challenges & Obstacles",
                    "content": [
                        "Technical limitations and barriers",
                        "Economic and financial constraints",
                        "Regulatory and compliance issues",
                        "Market adoption challenges",
                        "Competition and market dynamics"
                    ]
                },
                {
                    "title": "Innovations & Solutions",
                    "content": [
                        "Emerging technologies and approaches",
                        "Best practices and methodologies",
                        "Case studies and success stories",
                        "Investment and funding opportunities",
                        "Partnership and collaboration models"
                    ]
                },
                {
                    "title": "Future Outlook & Impact",
                    "content": [
                        "Projected growth and development",
                        "Potential market opportunities",
                        "Long-term implications and effects",
                        "Recommendations for stakeholders",
                        "Next steps and action items"
                    ]
                }
            ],
            "conclusion": {
                "title": "Conclusion",
                "content": [
                    f"{topic} represents a significant opportunity",
                    "Key challenges must be addressed strategically",
                    "Innovation and collaboration are essential",
                    "The future looks promising with proper execution"
                ]
            }
        }
    
    def create_presentation(self, content: Dict) -> Presentation:
        """
        Creates a PowerPoint presentation using python-pptx, now with images.
        """
        print("\n📊 Creating PowerPoint presentation...")
        prs = Presentation()
        
        # Add slides using the new helper methods
        self._add_title_slide(prs, content)
        self._add_content_slide(prs, content["overview"]["title"], content["overview"]["content"])

        for key_point in content["key_points"]:
            image_url = self.perform_image_search(key_point["title"])
            image_path = self.get_image_from_url(image_url) if image_url else None
            self._add_content_slide(prs, key_point["title"], key_point["content"], image_path)
        
        self._add_content_slide(prs, content["conclusion"]["title"], content["conclusion"]["content"])
        
        print("✅ PowerPoint presentation created successfully")
        return prs
    
    def save_presentation(self, prs: Presentation, topic: str) -> str:
        """
        Saves the presentation to a file with a dynamic filename.
        """
        safe_topic = re.sub(r'[^a-zA-Z0-9\s]', '', topic)
        safe_topic = re.sub(r'\s+', '_', safe_topic).lower()
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        filename = f"presentation_{safe_topic}_{timestamp}.pptx"
        
        try:
            prs.save(filename)
            print(f"💾 Presentation saved successfully: {filename}")
            return filename
        except Exception as e:
            print(f"❌ Error saving presentation: {e}")
            raise
    
    def generate_presentation(self, topic: Optional[str] = None) -> Optional[str]:
        """
        Main method to orchestrate the entire presentation generation process.
        Accepts a topic as an argument. If not provided, it prompts the user.
        """
        try:
            # Check for API keys at the start of generation
            if not self.gemini_api_key or not self.serpapi_api_key:
                raise ValueError("API keys are not configured. Please check your .env file.")

            if not topic:
                topic = self.get_user_topic()

            web_results = self.perform_web_search(topic)
            content = self.generate_llm_content(topic, web_results)
            presentation = self.create_presentation(content)
            filename = self.save_presentation(presentation, topic)
            
            print("\n" + "="*60)
            print("🎉 PRESENTATION GENERATION COMPLETE!")
            print("="*60)
            print(f"📁 File saved as: {filename}")
            print(f"📊 Total slides: {len(presentation.slides)}")
            print(f"🎯 Topic: {topic}")
            print("\nYou can now open the presentation in PowerPoint or any compatible application.")
            print("="*60)
            return filename
        except KeyboardInterrupt:
            print("\n\n⚠️  Operation cancelled by user.")
            return None
        except Exception as e:
            print(f"\n❌ An error occurred: {e}")
            print("Please check your API keys and try again.")
            return None

def main():
    """Main entry point of the script for command-line execution."""
    try:
        generator = PowerPointGenerator()
        generator.generate_presentation()
    except Exception as e:
        print(f"Fatal error: {e}")
        print("Please ensure all required libraries are installed and API keys are set correctly.")

if __name__ == "__main__":
    main()
